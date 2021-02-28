# -*- coding: utf-8 -*-
# Module: default
# License: GPL v.3 https://www.gnu.org/copyleft/gpl.html
from bs4 import BeautifulSoup
import requests
import os.path
from pptx import Presentation
from pptx.util import Inches, Pt

class ZonaFilms():
    
    
    @staticmethod    
    def get_parsed_page( url, headers = {}, params = None):
        headers = {
            "referer": "https://www.yandex.ru",
            "user-agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"
        }
        return BeautifulSoup(requests.get(url, headers = headers, params = params, timeout=(3.05, 27)).text, "lxml")
        
    

    def find_films(self):
        found_block = self.page.find("ul",{"class","results"})
        elements = found_block.find_all("li", {"class", "results-item-wrap"})
        #print(elements)
        found_items = []
        for el in elements:
            
            found_items.append({
            "pic":el.find("meta")["content"],
            "link":self.link + el.find("a")["href"],
            "name":el.find("a")["title"]
            })
        result = {}    
        for i in range(len(found_items)):    
            html = self.get_parsed_page(found_items[i]["link"]) 
            about_film = html.find("div", {"class", "entity-desc-table"})
            about_film = about_film.find_all("dd")
            film_name = found_items[i]["name"]
            film_img = found_items[i]["pic"]
            film_year = about_film[2].text.replace("\n", "")
            film_zhanr = about_film[3].text.replace("\n", "")
            film_country = about_film[4].text.replace("\n", "")
            film_director = about_film[5].text.replace("\n", "")
            film_actors = about_film[7].text.replace("\n", "").replace("  ","")
            film_description = html.find("div", {"class", "entity-desc-description"}).text
            full_info = {"film_img":film_img,"film_year":film_year,"film_zhanr":film_zhanr, "film_name":film_name,
                      "film_country":film_country,"film_director":film_director,"film_actors":film_actors,"film_description":film_description}
            result[i] = full_info
        return result
    def __init__(self,keyword, link = "https://w113.zona.plus"):
        self.link = link
        self.keyword  = keyword
        self.page = self.get_parsed_page(f"{link}/search/{keyword}")
        self.search_results = self.find_films()


class PresentationCreator():

    def __init__(self, film, ind = 0):
        self.film = film
        if os.path.isfile("result.pptx"):
            templ = "result.pptx"
        else:
            templ = "template.pptx"
        #___________________________________
        print(templ)
        self.prs = Presentation(templ)
        title_slide_layout = self.prs.slide_layouts[3]
        slide = self.prs.slides.add_slide(title_slide_layout)
        shapes = slide.shapes
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = self.film.search_results[ind]["film_name"]
        subtitle.text = "В главных ролях:\n"+self.film.search_results[ind]["film_actors"].replace(",", "\n")
        request_img = requests.get(self.film.search_results[ind]["film_img"]).content
        with open("img.jpg", "wb") as f:
            f.write(request_img)
        pic = slide.shapes.add_picture("img.jpg", Inches(.5), Inches(2.0))
        #___________________________________
        rows = 5
        cols = 2
        left = Inches(9.0)
        top = Inches(2.0)
        width = Inches(6.0)
        height = Inches(0.8)

        table = shapes.add_table(rows, cols, left, top, width, height).table


        table.columns[0].width = Inches(1.5)
        table.columns[1].width = Inches(2.5)


        table.cell(0, 0).text = 'Год'
        table.cell(1, 0).text = 'Страна'
        table.cell(2, 0).text = 'Слоган'
        table.cell(3, 0).text = 'Режисер'
        table.cell(4, 0).text = 'Жанр'


        table.cell(0, 1).text = self.film.search_results[ind]["film_year"]
        table.cell(1, 1).text = self.film.search_results[ind]["film_country"]
        table.cell(2, 1).text = ""
        table.cell(3, 1).text = self.film.search_results[ind]["film_director"]
        table.cell(4, 1).text = self.film.search_results[ind]["film_zhanr"]

        #_____________________________________

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = self.film.search_results[ind]["film_name"]
        subtitle.text = self.film.search_results[ind]["film_description"]

        #_____________________________________

        self.prs.save("result.pptx")
        



